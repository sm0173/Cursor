# -*- coding: utf-8 -*-
"""
Общая логика: PDF / DOCX / PPTX / (опционально PPT через LibreOffice) → Markdown, копирование *.sql.
Используется convert_bd_course_materials.py и convert_minor_course_materials.py.
"""
from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
from pathlib import Path


def _relative_md_basename(path: Path, src: Path) -> str:
    """Имя .md без расширения; для вложенных папок — через __."""
    path = path.resolve()
    src = src.resolve()
    rel = path.relative_to(src)
    parts = rel.with_suffix("").parts
    if len(parts) == 1:
        return parts[0]
    return "__".join(parts)


def pdf_to_markdown(pdf_path: Path) -> str:
    import fitz  # PyMuPDF

    parts: list[str] = [
        f"# {pdf_path.stem}\n",
        f"Источник (оригинал): `{pdf_path.name}`\n\n---\n\n",
    ]
    doc = fitz.open(pdf_path)
    try:
        for i in range(len(doc)):
            text = doc[i].get_text("text")
            if text.strip():
                parts.append(f"## Страница {i + 1}\n\n{text.rstrip()}\n\n")
    finally:
        doc.close()
    return "".join(parts).rstrip() + "\n"


def _heading_level_from_style(style_name: str) -> int | None:
    if not style_name:
        return None
    s = style_name.strip()
    low = s.lower()
    if low == "title":
        return 1
    m = re.search(r"(\d+)\s*$", s)
    if not m:
        return None
    n = int(m.group(1))
    if n < 1 or n > 6:
        return None
    if "heading" in low or "заголовок" in low or "rubrik" in low:
        return n
    if low.startswith("toc"):
        return None
    return None


def _escape_md_cell(text: str) -> str:
    return text.replace("\n", " ").replace("|", "\\|").strip()


def docx_to_markdown(docx_path: Path) -> str:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(docx_path)
    out: list[str] = [
        f"# {docx_path.stem}\n",
        f"Источник (оригинал): `{docx_path.name}`\n\n---\n\n",
    ]

    body = doc.element.body
    for child in body:
        tag = child.tag.split("}")[-1]
        if tag == "p":
            p = Paragraph(child, doc)
            t = p.text.strip()
            if not t:
                out.append("\n")
                continue
            lvl = _heading_level_from_style(p.style.name if p.style else "")
            if lvl is not None:
                out.append("#" * lvl + " " + t + "\n\n")
            else:
                out.append(t + "\n\n")
        elif tag == "tbl":
            table = Table(child, doc)
            rows_data: list[list[str]] = []
            for row in table.rows:
                rows_data.append([_escape_md_cell(c.text) for c in row.cells])
            if not rows_data:
                continue
            header = rows_data[0]
            out.append("| " + " | ".join(header) + " |\n")
            out.append("| " + " | ".join(["---"] * len(header)) + " |\n")
            for r in rows_data[1:]:
                while len(r) < len(header):
                    r.append("")
                out.append("| " + " | ".join(r[: len(header)]) + " |\n")
            out.append("\n")

    return "".join(out).rstrip() + "\n"


def _pptx_shape_texts(shape) -> list[str]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        t = shape.text_frame.text.strip()
        if t:
            out.append(t)
    if getattr(shape, "has_table", False) and shape.table:
        rows: list[list[str]] = []
        for row in shape.table.rows:
            rows.append([_escape_md_cell(c.text) for c in row.cells])
        if rows:
            hdr = rows[0]
            out.append("| " + " | ".join(hdr) + " |\n| " + " | ".join(["---"] * len(hdr)) + " |")
            for r in rows[1:]:
                while len(r) < len(hdr):
                    r.append("")
                out.append("| " + " | ".join(r[: len(hdr)]) + " |")
            out.append("")
    st = getattr(shape, "shape_type", None)
    if st == MSO_SHAPE_TYPE.GROUP and getattr(shape, "shapes", None):
        for child in shape.shapes:
            out.extend(_pptx_shape_texts(child))
    return out


def pptx_to_markdown(pptx_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    parts: list[str] = [
        f"# {pptx_path.stem}\n",
        f"Источник (оригинал): `{pptx_path.name}` (PowerPoint)\n\n---\n\n",
    ]
    for i, slide in enumerate(prs.slides, start=1):
        slide_chunks: list[str] = []
        for shape in slide.shapes:
            slide_chunks.extend(_pptx_shape_texts(shape))
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                nt = slide.notes_slide.notes_text_frame.text.strip()
                if nt:
                    slide_chunks.append("**Заметки к слайду:** " + nt)
        except (AttributeError, ValueError):
            pass
        body = "\n\n".join(s for s in slide_chunks if s).strip()
        if body:
            parts.append(f"## Слайд {i}\n\n{body}\n\n")
    return "".join(parts).rstrip() + "\n"


def _find_soffice() -> str | None:
    import shutil

    w = shutil.which("soffice")
    if w:
        return w
    for p in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if Path(p).is_file():
            return p
    return None


def ppt_to_markdown_via_libreoffice(ppt_path: Path) -> str | None:
    """
    Конвертация .ppt → PDF через LibreOffice, затем текст как у PDF.
    Возвращает None если soffice недоступен.
    """
    soffice = _find_soffice()
    if not soffice:
        return None
    tmp = Path(tempfile.mkdtemp(prefix="ppt_conv_"))
    try:
        cmd = [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp),
            str(ppt_path.resolve()),
        ]
        subprocess.run(cmd, check=True, timeout=180, capture_output=True)
        pdf_out = tmp / (ppt_path.stem + ".pdf")
        if not pdf_out.is_file():
            return None
        inner = pdf_to_markdown(pdf_out)
        sep = "\n---\n\n"
        pos = inner.find(sep)
        body = inner[pos + len(sep) :] if pos != -1 else inner
        return (
            f"# {ppt_path.stem}\n"
            f"Источник (оригинал): `{ppt_path.name}` (через LibreOffice → PDF → текст)\n\n"
            f"---\n\n"
            f"{body.strip()}\n"
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
        return None
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def ppt_stub_markdown(ppt_path: Path) -> str:
    return (
        f"# {ppt_path.stem}\n\n"
        f"Источник (оригинал): `{ppt_path.name}`\n\n"
        f"---\n\n"
        f"Автоматически извлечь текст из старого `.ppt` не получилось "
        f"(нужен **LibreOffice** в системе, команда `soffice`).\n\n"
        f"**Что сделать:** открыть файл в PowerPoint и «Сохранить как» PDF в ту же папку, "
        f"где лежат исходники, потом снова запустить конвертер — PDF обработается.\n"
    )


def convert_materials_folder(
    src: Path,
    dst: Path,
    *,
    readme_title: str,
    src_folder_name_for_text: str,
    run_command_hint: str,
) -> None:
    """
    src — папка с материалами (рекурсивно: pdf, docx, pptx, ppt, sql).
    dst — выходная папка с .md и копиями .sql.
    """
    src = src.resolve()
    if not src.is_dir():
        raise SystemExit(f"Не найдена папка-источник: {src}")

    dst.mkdir(parents=True, exist_ok=True)
    n_pdf = n_docx = n_pptx = n_ppt = n_sql = n_ppt_stub = 0

    for pdf in sorted(src.rglob("*.pdf")):
        if pdf.name.startswith("~$"):
            continue
        name = _relative_md_basename(pdf, src) + ".md"
        md = dst / name
        md.write_text(pdf_to_markdown(pdf), encoding="utf-8")
        print("OK PDF ->", md.name)
        n_pdf += 1

    for docx in sorted(src.rglob("*.docx")):
        if docx.name.startswith("~$"):
            continue
        name = _relative_md_basename(docx, src) + ".md"
        md = dst / name
        md.write_text(docx_to_markdown(docx), encoding="utf-8")
        print("OK DOCX ->", md.name)
        n_docx += 1

    for pptx in sorted(src.rglob("*.pptx")):
        if pptx.name.startswith("~$"):
            continue
        name = _relative_md_basename(pptx, src) + ".md"
        md = dst / name
        md.write_text(pptx_to_markdown(pptx), encoding="utf-8")
        print("OK PPTX ->", md.name)
        n_pptx += 1

    for ppt in sorted(src.rglob("*.ppt")):
        name = _relative_md_basename(ppt, src) + ".md"
        md = dst / name
        body = ppt_to_markdown_via_libreoffice(ppt)
        if body is None:
            md.write_text(ppt_stub_markdown(ppt), encoding="utf-8")
            print("STUB PPT (нет LibreOffice) ->", md.name)
            n_ppt_stub += 1
        else:
            md.write_text(body, encoding="utf-8")
            print("OK PPT ->", md.name)
        n_ppt += 1

    for sql in sorted(src.rglob("*.sql")):
        outp = dst / (_relative_md_basename(sql, src) + ".sql")
        shutil.copy2(sql, outp)
        print("OK SQL ->", outp.name)
        n_sql += 1

    readme = dst / "README.md"
    readme.write_text(
        f"""# {readme_title}

Здесь лежат выгрузки в **Markdown** (`.md`) из файлов папки `{src_folder_name_for_text}/`:
**PDF**, **DOCX**, **PPTX** (текст со слайдов; без OCR картинок), при наличии **LibreOffice** — **PPT** через промежуточный PDF, иначе для `.ppt` создаётся короткая заглушка с инструкцией.
Копируются также **.sql** (если есть).

Обход **рекурсивный** (`все подпапки`).

Пересобрать выгрузку:

```text
{run_command_hint}
```

Последний запуск: PDF — {n_pdf}, DOCX — {n_docx}, PPTX — {n_pptx}, PPT — {n_ppt} (из них заглушек без LO — {n_ppt_stub}), SQL — {n_sql}.
""",
        encoding="utf-8",
    )
    print("OK README.md")
